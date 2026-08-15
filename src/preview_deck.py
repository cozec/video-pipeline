"""Render a PNG preview of each deck slide by reading the .pptx back.

There is no LibreOffice on this machine, so this draws every shape at the position and
size stored in the file. Font metrics differ slightly from PowerPoint's, but it is enough
to see cramped spacing, misalignment and bad proportions before presenting.
"""

import matplotlib
matplotlib.use("Agg")

import matplotlib.patches as mpatches
import matplotlib.pyplot as plt
from PIL import Image
from pptx import Presentation
from pptx.enum.text import PP_ALIGN

from common import RESULTS, setup_logging

logger = setup_logging("preview")

FIGS = RESULTS / "slides"
DECK = FIGS / "video_pipeline_deck.pptx"
EMU = 914400.0


def rgb(color):
    try:
        return "#%02X%02X%02X" % tuple(color.rgb)
    except Exception:
        return "#000000"


def render() -> None:
    prs = Presentation(DECK)
    SW, SH = prs.slide_width / EMU, prs.slide_height / EMU

    for idx, slide in enumerate(prs.slides, 1):
        fig, ax = plt.subplots(figsize=(SW, SH), dpi=110)
        ax.set_xlim(0, SW); ax.set_ylim(SH, 0); ax.axis("off")
        ax.add_patch(mpatches.Rectangle((0, 0), SW, SH, facecolor="white", zorder=0))

        for sh in slide.shapes:
            if sh.left is None:
                continue
            x, y = sh.left / EMU, sh.top / EMU
            w, h = sh.width / EMU, sh.height / EMU

            if sh.shape_type == 13 or sh.__class__.__name__ == "Picture":
                try:
                    im = Image.open(FIGS / sh.image.filename) if getattr(sh.image, "filename", None) \
                         else Image.open(sh.image.blob and __import__("io").BytesIO(sh.image.blob))
                except Exception:
                    im = Image.open(__import__("io").BytesIO(sh.image.blob))
                ax.imshow(im, extent=(x, x + w, y + h, y), zorder=2, aspect="auto")
                continue

            if sh.has_text_frame and sh.text_frame.text.strip():
                # panels/rules carry no text; textboxes have no fill
                cy = y
                for p in sh.text_frame.paragraphs:
                    runs = [r for r in p.runs if r.text]
                    if not runs:
                        continue
                    r0 = runs[0]
                    pt = r0.font.size.pt if r0.font.size else 18
                    fam = "Menlo" if (r0.font.name or "") == "Menlo" else "Avenir Next"
                    col = rgb(r0.font.color) if r0.font.color and r0.font.color.type is not None else "#000000"
                    if p.space_before:
                        cy += p.space_before.pt / 72.0
                    txt = "".join(r.text for r in runs)
                    per = max(1, int(w / (pt * (0.60 if fam == "Menlo" else 0.55) / 72.0)))
                    lines, cur = [], ""
                    for word in txt.split(" "):
                        if len(cur) + len(word) + 1 <= per:
                            cur = (cur + " " + word).strip()
                        else:
                            lines.append(cur); cur = word
                    lines.append(cur)
                    lh = pt * (p.line_spacing or 1.15) / 72.0
                    ha = {PP_ALIGN.CENTER: "center", PP_ALIGN.RIGHT: "right"}.get(p.alignment, "left")
                    tx = x + (w / 2 if ha == "center" else (w if ha == "right" else 0))
                    for ln in lines:
                        cy += lh
                        ax.text(tx, cy - lh * 0.24, ln, fontsize=pt * 0.98, color=col,
                                family=fam, ha=ha, va="baseline", zorder=4,
                                fontweight="bold" if r0.font.bold else "normal")
                continue

            # non-text shapes: panels and rules
            fill = "none"
            try:
                if sh.fill.type is not None and sh.fill.type == 1:
                    fill = rgb(sh.fill.fore_color)
            except Exception:
                pass
            edge = "none"
            try:
                if sh.line.color and sh.line.color.type is not None:
                    edge = rgb(sh.line.color)
            except Exception:
                pass
            ax.add_patch(mpatches.FancyBboxPatch(
                (x, y), w, max(h, 0.008),
                boxstyle="round,pad=0,rounding_size=0.06",
                facecolor=fill, edgecolor=edge, linewidth=0.8, zorder=1))

        out = FIGS / f"preview_slide{idx}.png"
        fig.savefig(out, dpi=110, bbox_inches="tight", pad_inches=0)
        plt.close(fig)
        logger.info("wrote %s", out.name)


if __name__ == "__main__":
    render()
